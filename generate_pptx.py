from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16.54)   # A3 landscape
prs.slide_height = Inches(11.69)

# Colors
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE   = RGBColor(0xE6, 0x7E, 0x22)
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)
C_BLUE     = RGBColor(0x29, 0x80, 0xB9)
C_PURPLE   = RGBColor(0x8E, 0x44, 0xAD)
C_DARK     = RGBColor(0x22, 0x22, 0x22)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xFA, 0xFA, 0xFA)
C_GREY_BG  = RGBColor(0xEE, 0xEE, 0xEE)
C_GREY_TXT = RGBColor(0x88, 0x88, 0x88)
C_TAG_BG   = RGBColor(0x34, 0x49, 0x5E)
C_DANGER_BG = RGBColor(0xFD, 0xED, 0xEC)
C_ACTION_BG = RGBColor(0xEA, 0xFA, 0xF1)
C_AUTO_BG   = RGBColor(0xEA, 0xF2, 0xFF)
C_HUMAN_BG  = RGBColor(0xFD, 0xF2, 0xE9)
C_SMS_BG    = RGBColor(0xF4, 0xEC, 0xF7)
C_DECISION  = RGBColor(0xFF, 0xF3, 0xCD)

def add_textbox(slide, left, top, width, height, text, font_size=9,
                bold=False, color=C_DARK, bg=None, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    if anchor:
        tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Meiryo"
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    return txBox

def add_rect(slide, left, top, width, height, fill, text="",
             font_size=10, font_color=C_WHITE, bold=True, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Meiryo"
    return shape

def add_bordered_box(slide, left, top, width, height, fill, border_color,
                     text="", font_size=8, font_color=C_DARK, bold=False,
                     align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Meiryo"
    return shape

def add_diamond(slide, left, top, width, height, text, fill=C_DECISION):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C_ORANGE
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.name = "Meiryo"
    return shape

# ── Layout constants ──
LEFT_MARGIN = 0.2
COL_PHASE_W = 1.1
LANE_W = 2.85
GAP = 0.04
lane_x = [LEFT_MARGIN + COL_PHASE_W + GAP + i * (LANE_W + GAP) for i in range(5)]
phase_x = LEFT_MARGIN

# Row heights and Y positions
HEADER_H = 0.35
LANE_HDR_H = 0.30
row_y_start = 1.50
ROW_H = [0.95, 1.05, 2.85, 2.45, 2.50, 0.55]
row_y = []
y = row_y_start
for h in ROW_H:
    row_y.append(y)
    y += h + GAP

# ── Slide 1 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

# Title
add_textbox(slide, 0, 0.1, 16.54, 0.5,
            "Fraud Kill Chain（不正送金キルチェーン）",
            font_size=20, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.55, 16.54, 0.3,
            'シート名「不審な電話手口」 ― 犯人サイド × 被害企業サイド スイムレーン図',
            font_size=10, color=C_GREY_TXT, align=PP_ALIGN.CENTER)

# ── Key info banners ──
add_rect(slide, 0.5, 0.88, 3.2, 0.30,
         C_RED, "対象：法人向けIB（BizSTATION）", 10, C_WHITE)
add_rect(slide, 3.85, 0.88, 6.8, 0.30,
         C_DARK, "犯人はBizSTATIONの認証フロー（契約者番号→利用者ID→ログインPW→取引実行PW→OTP）を完全に把握", 8, C_WHITE, bold=False)
add_rect(slide, 10.80, 0.88, 5.2, 0.30,
         C_ORANGE, "法人IBは送金限度額が高く、1件で数千万〜数億円の被害リスク", 8, C_WHITE, bold=False)

# ── Actor group headers ──
actor_y = 1.25
add_rect(slide, lane_x[0], actor_y, LANE_W * 2 + GAP, HEADER_H,
         C_DARK, "犯人サイド", 11)
add_rect(slide, lane_x[2], actor_y, LANE_W * 3 + GAP * 2, HEADER_H,
         C_ORANGE, "被害企業サイド", 11)

# ── Lane headers ──
lh_y = actor_y + HEADER_H + GAP
lane_names = ["情報・インフラ", "架電", "受電", "会社のPC", "個人スマホ"]
add_rect(slide, phase_x, lh_y, COL_PHASE_W, LANE_HDR_H,
         C_GREY_BG, "フェーズ", 9, C_DARK)
for i, name in enumerate(lane_names):
    add_rect(slide, lane_x[i], lh_y, LANE_W, LANE_HDR_H,
             C_GREY_BG, name, 9, C_DARK)

# ── Phase labels ──
phase_labels = [
    "① 犯行準備\n（偵察など）",
    "② 被害者への\n接触",
    "③ 心理的操作",
    "④ システムや\n機能の侵害",
    "⑤ 侵害に使用\nする認証情報",
    "⑥ 収益化\n／マネロン",
]
for i, label in enumerate(phase_labels):
    add_rect(slide, phase_x, row_y[i], COL_PHASE_W, ROW_H[i],
             C_RED, label, 10, C_WHITE)

# ── Empty cell backgrounds ──
for i in range(6):
    for j in range(5):
        add_rect(slide, lane_x[j], row_y[i], LANE_W, ROW_H[i],
                 C_WHITE, "", font_size=1, font_color=C_WHITE)

# ── Phase 1 content ──
r = 0
add_bordered_box(slide, lane_x[0] + 0.1, row_y[r] + 0.15, LANE_W - 0.2, 0.55,
                 C_ACTION_BG, C_GREEN,
                 "電話番号・メールアドレスを取得", 8, bold=True)
add_bordered_box(slide, lane_x[1] + 0.1, row_y[r] + 0.15, LANE_W - 0.2, 0.35,
                 C_HUMAN_BG, C_ORANGE, "事前調査", 8)
add_textbox(slide, lane_x[1] + 0.1, row_y[r] + 0.55, LANE_W - 0.2, 0.3,
            "※幅広に架電している印象", 7, color=C_GREY_TXT)

# ── Phase 2 content ──
r = 1
add_bordered_box(slide, lane_x[1] + 0.1, row_y[r] + 0.08, LANE_W - 0.2, 0.85,
                 C_AUTO_BG, C_BLUE,
                 "【自動音声】\n「こちらはMバンく BizSTATION法人窓口です。"
                 "BizSTATIONを利用している方は1を押してください。"
                 "承認実行権限を持っている方は1を押してください」", 7)
add_bordered_box(slide, lane_x[2] + 0.1, row_y[r] + 0.08, LANE_W - 0.2, 0.35,
                 C_ACTION_BG, C_GREEN, "「1」押下", 9, bold=True)
add_textbox(slide, lane_x[2] + 0.1, row_y[r] + 0.50, LANE_W - 0.2, 0.55,
            "※「1」以外を押下した人は、法人ダイレクトオフィス（コールセンター）の電話番号を案内される",
            6, color=C_GREY_TXT)

# ── Phase 3 content ──
r = 2
cy = row_y[r]
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.05, LANE_W - 0.2, 0.70,
                 C_HUMAN_BG, C_ORANGE,
                 "【有人対応】\n「パソコン環境の更新が必要。現在口座の利用制限をしています。"
                 "手続きを郵送していますが、届いてますか。更新されていないので電話しました」",
                 7)
add_textbox(slide, lane_x[1] + 0.1, cy + 0.78, LANE_W - 0.2, 0.2,
            "※複数の誘導方法あり", 6, color=C_GREY_TXT)
add_textbox(slide, lane_x[1] + 0.1, cy + 0.90, LANE_W - 0.2, 0.2,
            "※AI音声または有人対応（法人高額案件のため両方の可能性あり）", 6, bold=True, color=C_RED)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 1.02, LANE_W - 0.2, 0.22,
                 C_DARK, C_DARK,
                 "🔑 攻撃の鍵：「口座制限」で恐怖→以降の判断力を奪う", 6, C_WHITE, bold=True)

add_diamond(slide, lane_x[1] + 0.7, cy + 1.28, 1.4, 0.50,
            "メアド\n入手済／未済？")

add_bordered_box(slide, lane_x[1] + 0.1, cy + 1.82, LANE_W - 0.2, 0.40,
                 C_HUMAN_BG, C_ORANGE,
                 "【未済の場合】\n「手続きをメールで案内するので、メールアドレスを教えてください」",
                 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 2.24, LANE_W - 0.2, 0.30,
                 C_DANGER_BG, C_RED,
                 "⚠ 破綻ポイント：BizSTATION契約時にメアド登録済み。銀行が聞く必要はない→手順②の恐怖で気づけない", 6)

add_bordered_box(slide, lane_x[1] + 0.1, cy + 2.56, LANE_W - 0.2, 0.22,
                 C_HUMAN_BG, C_ORANGE,
                 "→ メール送付", 8, bold=True)

add_bordered_box(slide, lane_x[1] + 0.1, cy + 2.80, LANE_W - 0.2, 0.15,
                 C_LIGHT, C_GREY_TXT,
                 '件名①「【重要】BizSTATIONお客さま情報更新のお知らせ」', 6)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 2.97, LANE_W - 0.2, 0.15,
                 C_LIGHT, C_GREY_TXT,
                 '件名②「【Mバンく】IB安全利用案内（Rapport推奨）」', 6)

# Victim side phase 3
add_bordered_box(slide, lane_x[2] + 0.1, cy + 1.82, LANE_W - 0.2, 0.30,
                 C_ACTION_BG, C_GREEN,
                 "メールアドレスを伝える", 8)
add_bordered_box(slide, lane_x[3] + 0.1, cy + 2.56, LANE_W - 0.2, 0.35,
                 C_ACTION_BG, C_GREEN,
                 "メール受信\n（Bizデザインのメール）", 8)

# ── Phase 4 content ──
r = 3
cy = row_y[r]
# Attacker calls
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.05, LANE_W - 0.2, 0.50,
                 C_HUMAN_BG, C_ORANGE,
                 "【有人対応】\n「セキュリティ強化のためソフトインストールが必要です。"
                 "メールのボタンをクリックしてください」", 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.60, LANE_W - 0.2, 0.30,
                 C_HUMAN_BG, C_ORANGE,
                 "「セットアップをクリックしてください」", 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.95, LANE_W - 0.2, 0.30,
                 C_HUMAN_BG, C_ORANGE,
                 "「実行するをクリックしてください」", 7)

# Company PC
add_bordered_box(slide, lane_x[3] + 0.1, cy + 0.05, LANE_W - 0.2, 0.35,
                 C_ACTION_BG, C_GREEN,
                 "クリック → ブラウザ表示\n（Bizデザイン重要なお知らせ画面）", 7)
add_bordered_box(slide, lane_x[3] + 0.1, cy + 0.45, LANE_W - 0.2, 0.35,
                 C_DANGER_BG, C_RED,
                 "【DL】リモートツールのセットアップファイル setup.msi がダウンロードされる", 7)
add_bordered_box(slide, lane_x[3] + 0.1, cy + 0.85, LANE_W - 0.2, 0.40,
                 C_DANGER_BG, C_RED,
                 "【実行】msiファイル実行 → ポップアップ「実行する/実行しない」\n※「不明な発行元」と表示", 7)
add_bordered_box(slide, lane_x[3] + 0.1, cy + 1.30, LANE_W - 0.2, 0.30,
                 C_DANGER_BG, C_RED,
                 "【インストール】ScreenConnect（リモートツール）がインストールされる", 7, bold=True)

# Red screen
add_bordered_box(slide, lane_x[3] + 0.1, cy + 1.65, LANE_W - 0.2, 0.85,
                 C_RED, C_RED,
                 "『Mバンく　現在セキュリティ保護強化に伴うシステム更新ならびに、"
                 "お客さま情報の確認、更新処理を実施しております。"
                 "処理完了までの間一時的に画面が遅延または停止する場合がありますが、"
                 "正常な処理となりますのでそのままお待ちくださいますようお願い申し上げます。』",
                 7, C_WHITE)

# ── Phase 5 content ──
r = 4
cy = row_y[r]
# Attacker calls
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.05, LANE_W - 0.2, 0.40,
                 C_HUMAN_BG, C_ORANGE,
                 "「手続きを進めるにはSMS認証が必要なため、お客さまの携帯電話の番号を教えてください」", 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.50, LANE_W - 0.2, 0.22,
                 C_SMS_BG, C_PURPLE, "【SMS送付】", 8, bold=True)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.78, LANE_W - 0.2, 0.30,
                 C_HUMAN_BG, C_ORANGE,
                 "「認証画面を表示するためリンクをクリックしてください」", 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 1.13, LANE_W - 0.2, 0.45,
                 C_HUMAN_BG, C_ORANGE,
                 "「認証のため契約者番号、利用者ID、ログインPW、取引実行PW、OTPを入力してください」", 7)
add_bordered_box(slide, lane_x[1] + 0.1, cy + 1.65, LANE_W - 0.2, 0.45,
                 C_DANGER_BG, C_RED,
                 "【犯人ブラウザ】赤い画面の裏でBizにアクセス、窃取したOTPを入力", 7, bold=True)

# Victim: phone number
add_bordered_box(slide, lane_x[2] + 0.1, cy + 0.05, LANE_W - 0.2, 0.30,
                 C_ACTION_BG, C_GREEN, "電話番号を伝える", 8)

# Victim: smartphone
add_bordered_box(slide, lane_x[4] + 0.1, cy + 0.50, LANE_W - 0.2, 0.30,
                 C_SMS_BG, C_PURPLE, "SMS受信 → リンクをクリック", 8)
add_bordered_box(slide, lane_x[4] + 0.1, cy + 0.88, LANE_W - 0.2, 0.70,
                 C_DANGER_BG, C_RED,
                 "開いた画面に認証情報を入力\n"
                 "・契約者番号\n・利用者ID\n・ログインPW\n・取引実行PW\n・OTP", 7)

# ── Phase 6 content ──
r = 5
cy = row_y[r]
add_bordered_box(slide, lane_x[1] + 0.1, cy + 0.08, LANE_W - 0.2, 0.45,
                 C_RED, C_RED,
                 "不正送金", 14, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── Footer ──
add_textbox(slide, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 不審な電話手口 ｜ A3横印刷対応",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 2: 銀行の正規フローとの類似性
# ══════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide2, 0, 0.3, 16.54, 0.6,
            "なぜ見破れないのか — 銀行の正規フローとの類似性",
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_RED)
add_textbox(slide2, 0, 0.85, 16.54, 0.4,
            "銀行自身が顧客を「自動音声の指示に従う」行動パターンに教育してきた結果、この攻撃が成立している。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# Comparison table
from pptx.util import Inches, Pt, Emu
tbl_left = 1.5
tbl_top = 1.6
tbl_w = 13.5
col_w = tbl_w / 3
row_h_tbl = 0.55
headers = ["銀行の正規フロー", "この攻撃の手口", "顧客の体験"]
rows = [
    ["自動音声で用件を振り分け", "自動音声で権限者を選別", "同じ"],
    ["番号を押して選択", "「1」を押す", "同じ"],
    ["担当者と話す", "犯人（AI音声/有人）と話す", "同じ"],
    ["指示に従って手続き", "指示に従ってインストール", "同じ"],
    ["本人確認で情報を伝える", "認証情報をすべて入力", "同じ"],
]

# Header row
for ci, h in enumerate(headers):
    add_rect(slide2, tbl_left + ci * col_w, tbl_top, col_w, row_h_tbl,
             C_DARK, h, 12, C_WHITE)

# Data rows
for ri, row in enumerate(rows):
    y = tbl_top + (ri + 1) * row_h_tbl
    bg_even = RGBColor(0xFF, 0xFF, 0xFF)
    bg_odd = RGBColor(0xF9, 0xF9, 0xF9)
    bg = bg_odd if ri % 2 else bg_even
    for ci, cell in enumerate(row):
        if ci == 2:
            add_rect(slide2, tbl_left + ci * col_w, y, col_w, row_h_tbl,
                     bg, cell, 14, C_RED, bold=True)
        else:
            add_bordered_box(slide2, tbl_left + ci * col_w, y, col_w, row_h_tbl,
                             bg, RGBColor(0xCC, 0xCC, 0xCC), cell, 11, C_DARK)

# Conclusion
add_rect(slide2, 1.5, tbl_top + 6 * row_h_tbl + 0.3, 13.5, 0.50,
         C_RED, "構造的問題：正規と詐欺の区別が顧客側ではほぼ不可能", 16, C_WHITE)

add_textbox(slide2, 1.5, tbl_top + 6 * row_h_tbl + 1.1, 13.5, 1.8,
            "銀行側が取るべき対策：\n"
            "① 架電を原則廃止し「銀行からは電話しません」と宣言する\n"
            "② やむを得ない架電は事前にIB画面やアプリで通知してから行う\n"
            "③「電話でソフトインストールを指示することはない」を明示する（現状周知不足）",
            font_size=13, color=C_DARK)

add_textbox(slide2, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 銀行の正規フローとの類似性",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 3: 手順②の心理操作テクニック
# ══════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide3, 0, 0.3, 16.54, 0.6,
            "手順②の心理操作 — どう答えても犯人のシナリオに乗る設計",
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_ORANGE)
add_textbox(slide3, 0, 0.85, 16.54, 0.4,
            "4つのセリフそれぞれに明確な狙いがあり、被害者の心理的退路を塞いでいる。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# Technique table
s3_top = 1.5
s3_cols = [4.5, 2.0, 7.0]
s3_left = 1.5
s3_row_h = 0.65
s3_headers = ["セリフ", "狙い", "効果"]
s3_rows = [
    ["「パソコン環境の更新が必要」", "伏線", "後のソフトインストール（手順⑤⑥）を自然に受け入れさせる"],
    ["「口座の利用制限をしています」", "恐怖", "「今すぐ対応しなきゃ」と焦らせ、冷静な判断を奪う"],
    ["「手続きを郵送しましたが、届いてますか」", "負い目", "被害者に「自分が見落とした」と思わせ、指示に従いやすくする"],
    ["「更新されていないので電話しました」", "正当化", "架電がフォローアップに見え、銀行の丁寧な対応と錯覚させる"],
]
s3_tag_colors = [C_BLUE, C_RED, C_ORANGE, C_PURPLE]

# Header
x_offset = s3_left
for ci, h in enumerate(s3_headers):
    add_rect(slide3, x_offset, s3_top, s3_cols[ci], s3_row_h,
             C_DARK, h, 12, C_WHITE)
    x_offset += s3_cols[ci]

# Data
for ri, row in enumerate(s3_rows):
    y = s3_top + (ri + 1) * s3_row_h
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x_offset = s3_left
    for ci, cell in enumerate(row):
        if ci == 1:
            add_rect(slide3, x_offset, y, s3_cols[ci], s3_row_h,
                     s3_tag_colors[ri], cell, 12, C_WHITE, bold=True)
        else:
            add_bordered_box(slide3, x_offset, y, s3_cols[ci], s3_row_h,
                             bg, RGBColor(0xCC, 0xCC, 0xCC), cell, 10, C_DARK)
        x_offset += s3_cols[ci]

# Response pattern section
rp_y = s3_top + 5 * s3_row_h + 0.3
add_rect(slide3, s3_left, rp_y, 13.5, 0.45,
         C_DECISION, "「届いてますか？」への回答パターン", 13, C_DARK, bold=True, align=PP_ALIGN.LEFT)

add_bordered_box(slide3, s3_left, rp_y + 0.50, 6.5, 0.50,
                 RGBColor(0xFF, 0xFF, 0xFF), C_ORANGE,
                 "「届いてない」（大多数）\n→「届いていないなら仕方ない、電話で手続きしましょう」", 10)
add_rect(slide3, s3_left + 6.6, rp_y + 0.50, 1.5, 0.50,
         C_RED, "→ シナリオ続行", 10, C_WHITE)

add_bordered_box(slide3, s3_left + 8.3, rp_y + 0.50, 3.8, 0.50,
                 RGBColor(0xFF, 0xFF, 0xFF), C_ORANGE,
                 "「届いてる」（稀）\n→「では確認しながら進めましょう」", 10)
add_rect(slide3, s3_left + 12.2, rp_y + 0.50, 1.3, 0.50,
         C_RED, "→ シナリオ続行", 10, C_WHITE)

# Conclusion
add_rect(slide3, s3_left, rp_y + 1.2, 13.5, 0.50,
         C_RED, "どう答えても犯人のシナリオから逸脱できない — 顧客には見破れない", 15, C_WHITE)

add_textbox(slide3, s3_left, rp_y + 1.9, 13.5, 0.6,
            "これも銀行が日常的に行う「書類を郵送しました → 届きましたか → 未対応なのでお電話」と同一パターンであり、顧客には見破れない",
            font_size=12, color=C_DARK)

add_textbox(slide3, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 手順②の心理操作テクニック",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 4: 顧客は見破れるか — 結論
# ══════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide4, 0, 0.2, 16.54, 0.6,
            "顧客はこの攻撃を見破れるか？",
            font_size=24, bold=True, align=PP_ALIGN.CENTER, color=C_RED)

# Detection table
s4_left = 0.8
s4_top = 1.0
s4_cols = [1.0, 3.8, 1.8, 6.5]
s4_row_h = 0.50
s4_headers = ["手順", "内容", "見破れるか", "理由"]

# Header
x = s4_left
for ci, h in enumerate(s4_headers):
    add_rect(slide4, x, s4_top, s4_cols[ci], s4_row_h,
             C_DARK, h, 10, C_WHITE)
    x += s4_cols[ci]

C_IMPOSSIBLE = RGBColor(0xC0, 0x39, 0x2B)
C_DIFFICULT  = RGBColor(0xE6, 0x7E, 0x22)
C_CHANCE     = RGBColor(0x27, 0xAE, 0x60)

s4_rows = [
    ["①", "自動音声「1を押してください」",       "不可能",       "銀行も同じことをする",                C_IMPOSSIBLE],
    ["②", "「口座制限中、郵送したが届いてますか」", "不可能",       "銀行も同じことを言う",                C_IMPOSSIBLE],
    ["③", "「メールで案内します」メアド聞き出し",   "ほぼ不可能",   "銀行も確認することがある",              C_IMPOSSIBLE],
    ["④", "Bizデザインの偽メール送付",            "困難",         "精巧に模倣されている",                 C_DIFFICULT],
    ["⑤", "「ソフトをインストールしてください」",   "唯一のチャンス", "ただし既に心理的に従う状態",           C_CHANCE],
    ["⑥", "msi実行「不明な発行元」の警告",        "技術的には可能", "電話で「実行して」と言われ突破される",    C_DIFFICULT],
    ["⑦", "赤い画面「お待ちください」",           "困難",         "もう完全に信用している",               C_DIFFICULT],
    ["⑧", "認証情報5要素を一画面で入力",          "本来は異常",    "焦りと電話指示で気づく余裕がない",      C_DIFFICULT],
]

for ri, row in enumerate(s4_rows):
    y = s4_top + (ri + 1) * s4_row_h
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x = s4_left
    # 手順
    add_rect(slide4, x, y, s4_cols[0], s4_row_h,
             bg, row[0], 11, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    x += s4_cols[0]
    # 内容
    add_bordered_box(slide4, x, y, s4_cols[1], s4_row_h,
                     bg, RGBColor(0xDD, 0xDD, 0xDD), row[1], 9, C_DARK)
    x += s4_cols[1]
    # 見破れるか
    add_rect(slide4, x, y, s4_cols[2], s4_row_h,
             row[4], row[2], 9, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += s4_cols[2]
    # 理由
    add_bordered_box(slide4, x, y, s4_cols[3], s4_row_h,
                     bg, RGBColor(0xDD, 0xDD, 0xDD), row[3], 9, C_DARK)

# Arrow pointing to the key insight
key_y = s4_top + 9 * s4_row_h + 0.25
add_rect(slide4, s4_left, key_y, 13.1, 0.55,
         C_RED, "唯一の見破りポイント（手順⑤⑥）に到達した時点で、冷静に判断できる心理状態ではない", 14, C_WHITE)

# Two-column conclusion
conc_y = key_y + 0.75
# Left: Customer (powerless)
add_rect(slide4, s4_left, conc_y, 6.4, 0.40,
         C_GREY_BG, "顧客側（ほぼ無力）", 12, C_DARK, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide4, s4_left, conc_y + 0.40, 6.4, 0.50,
                 RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xDD, 0xDD, 0xDD),
                 "「気をつける」「注意喚起」では構造的に防げない", 10, C_GREY_TXT)

# Right: Bank/Enterprise (only defense)
add_rect(slide4, s4_left + 6.7, conc_y, 6.4, 0.40,
         C_RED, "銀行/企業側（ここでしか止められない）", 12, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide4, s4_left + 6.7, conc_y + 0.40, 6.4, 1.60,
                 RGBColor(0xFF, 0xFF, 0xFF), C_RED,
                 "① ScreenConnectの通信をNWで検知\n"
                 "② 送金時の振る舞い検知（時間帯、金額、送金先）\n"
                 "③ 架電を廃止し「電話では手続きしない」と宣言\n"
                 "④ リモートツール実行をEDRでブロック\n"
                 "⑤ SMS認証のリンク先ドメインをホワイトリスト化",
                 10, C_DARK)

# Final verdict
add_rect(slide4, s4_left, conc_y + 2.2, 13.1, 0.50,
         C_DARK, "結論：「顧客が気をつければ防げる」は幻想 — 銀行側のシステム対策でしか止められない", 14, C_WHITE)

add_textbox(slide4, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 顧客は見破れるか",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 5: 1本の電話 — つないだまま全手順を完了させる
# ══════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide5.background.fill.solid()
slide5.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide5, 0, 0.2, 16.54, 0.6,
            "1本の電話 — つないだまま全手順を完了させる",
            font_size=24, bold=True, align=PP_ALIGN.CENTER, color=C_RED)
add_textbox(slide5, 0, 0.75, 16.54, 0.4,
            "この攻撃は「9つの手順」ではなく、1本の電話の中で起きる1つの連続した攻撃である。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# Timeline - vertical flow
tl_left = 1.0
tl_step_w = 5.5
tl_bar_x = tl_left + tl_step_w + 0.2
tl_bar_w = 0.5
tl_top = 1.3
tl_step_h = 0.42
tl_gap = 0.03

steps = [
    ("①", "自動音声で架電",              "電話接続"),
    ("②", "心理操作（口座制限中）",       ""),
    ("③", "メアド聞き出し＋メール送付",    ""),
    ("④", "「メールを開いてください」",     ""),
    ("⑤", "「インストールしてください」",   "ずっと"),
    ("⑥", "「実行してください」",          "電話が"),
    ("⑦", "赤い画面表示",                "つながっている"),
    ("⑧", "「携帯番号を教えて」",          ""),
    ("⑨", "「認証情報を入力して」",        ""),
    ("",   "不正送金",                    "電話終了？"),
]

for i, (num, desc, bar_text) in enumerate(steps):
    y = tl_top + i * (tl_step_h + tl_gap)
    if num == "":
        add_rect(slide5, tl_left, y, tl_step_w, tl_step_h,
                 C_RED, f"  {desc}", 11, C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    else:
        bg = RGBColor(0xFF, 0xFF, 0xFF)
        add_bordered_box(slide5, tl_left, y, tl_step_w, tl_step_h,
                         bg, RGBColor(0xCC, 0xCC, 0xCC),
                         f"{num}  {desc}", 10, C_DARK)

# Continuous phone bar
bar_total_h = 10 * (tl_step_h + tl_gap) - tl_gap
add_rect(slide5, tl_bar_x, tl_top, tl_bar_w, bar_total_h,
         C_RED, "", 1, C_WHITE)
# Bar labels
add_textbox(slide5, tl_bar_x - 0.1, tl_top + 0.1, tl_bar_w + 0.2, 0.3,
            "電話接続", 8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide5, tl_bar_x - 0.3, tl_top + bar_total_h * 0.4, tl_bar_w + 0.6, 0.5,
            "ずっと\n電話が\nつながっている", 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Right side: effects table
eff_left = 8.0
eff_top = 1.3
eff_w = 7.5
eff_row_h = 0.55

add_rect(slide5, eff_left, eff_top, eff_w, eff_row_h,
         C_DARK, "電話をつないだままにする5つの効果", 13, C_WHITE)

effects = [
    ("考える隙を与えない",     "電話を切ったら冷静になる。切らせない", C_RED),
    ("リアルタイム制御",        "被害者の画面状況に合わせて次の指示を出せる", C_ORANGE),
    ("離脱防止",              "「このままお待ちください」で引き留められる", C_BLUE),
    ("第三者への相談を防ぐ",    "電話中なので同僚や家族に相談できない", C_PURPLE),
    ("信頼の維持",             "銀行員が丁寧に付き添ってくれている感覚", C_GREEN),
]

for i, (title, desc, color) in enumerate(effects):
    y = eff_top + (i + 1) * eff_row_h
    add_rect(slide5, eff_left, y, 2.5, eff_row_h,
             color, title, 10, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    bg = RGBColor(0xF9, 0xF9, 0xF9) if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    add_bordered_box(slide5, eff_left + 2.5, y, eff_w - 2.5, eff_row_h,
                     bg, RGBColor(0xDD, 0xDD, 0xDD), desc, 10, C_DARK)

# "If the call was disconnected" section
disc_y = eff_top + 6 * eff_row_h + 0.25
add_rect(slide5, eff_left, disc_y, eff_w, 0.40,
         C_GREEN, "もし電話が切れたら（攻撃失敗）", 11, C_WHITE)
disc_items = [
    "「本当に銀行からの電話だったのか？」と疑問に思う",
    "銀行の公式番号にかけ直す",
    "同僚や上司に相談する",
    "ネットで調べる",
]
add_bordered_box(slide5, eff_left, disc_y + 0.40, eff_w, 0.90,
                 RGBColor(0xFF, 0xFF, 0xFF), C_GREEN,
                 "\n".join(f"→ {item}" for item in disc_items),
                 10, C_DARK)

# Bottom conclusion
add_rect(slide5, 1.0, disc_y + 1.5, 14.5, 0.50,
         C_DARK,
         "被害者の体験：「銀行の担当者に電話で付き添ってもらいながら手続きをした」— それが攻撃だったとは気づかない",
         13, C_WHITE)

add_textbox(slide5, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 1本の電話の構造",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 6: 手順④ フィッシングメール — 迷惑メールフィルタも無力化
# ══════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
slide6.background.fill.solid()
slide6.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide6, 0, 0.2, 16.54, 0.6,
            "手順④ フィッシングメール — セキュリティ対策も電話1本で無力化",
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_RED)
add_textbox(slide6, 0, 0.75, 16.54, 0.4,
            "電話で予告してからメールを送る。通常のフィッシングとは根本的に異なる構造。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# Two-pattern comparison
s6_left = 1.0
s6_top = 1.3

# Left: normal phishing
add_rect(slide6, s6_left, s6_top, 7.0, 0.45,
         C_GREY_BG, "通常のフィッシング", 14, C_DARK, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left, s6_top + 0.50, 7.0, 0.45,
                 RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xCC, 0xCC, 0xCC),
                 "メールが突然届く", 12, C_DARK)
add_rect(slide6, s6_left + 1.5, s6_top + 1.00, 4.0, 0.40,
         C_GREY_BG, "↓", 14, C_DARK, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left, s6_top + 1.45, 7.0, 0.45,
                 RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xCC, 0xCC, 0xCC),
                 "「怪しいかも？」と思う余地あり", 12, C_DARK)
add_rect(slide6, s6_left + 1.5, s6_top + 1.95, 4.0, 0.40,
         C_GREY_BG, "↓", 14, C_DARK, align=PP_ALIGN.CENTER)
add_rect(slide6, s6_left, s6_top + 2.40, 7.0, 0.45,
         C_GREEN, "迷惑メールフィルタが防いでくれる可能性あり", 11, C_WHITE)

# Right: this attack
add_rect(slide6, s6_left + 7.5, s6_top, 7.0, 0.45,
         C_RED, "この攻撃", 14, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left + 7.5, s6_top + 0.50, 7.0, 0.45,
                 RGBColor(0xFF, 0xFF, 0xFF), C_RED,
                 "電話で「今からメール送ります」と予告", 12, C_DARK, bold=True)
add_rect(slide6, s6_left + 9.0, s6_top + 1.00, 4.0, 0.40,
         RGBColor(0xFD, 0xED, 0xEC), "↓", 14, C_RED, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left + 7.5, s6_top + 1.45, 7.0, 0.45,
                 RGBColor(0xFF, 0xFF, 0xFF), C_RED,
                 "メール到着 →「本物だ」と確信（予告通り来たから）", 12, C_DARK, bold=True)
add_rect(slide6, s6_left + 9.0, s6_top + 1.95, 4.0, 0.40,
         RGBColor(0xFD, 0xED, 0xEC), "↓", 14, C_RED, align=PP_ALIGN.CENTER)
add_rect(slide6, s6_left + 7.5, s6_top + 2.40, 7.0, 0.45,
         C_RED, "電話 × メールの二重チャネルが相互に信頼性を補強", 11, C_WHITE)

# Spam filter bypass section
sf_y = s6_top + 3.2
add_rect(slide6, s6_left, sf_y, 14.5, 0.50,
         C_DARK, "偽メールは正規フォルダに入らない（SPF/DKIM/DMARC不通過）→ ほぼ確実に迷惑フォルダ行き", 12, C_WHITE)

# Flow: spam scenario
add_bordered_box(slide6, s6_left, sf_y + 0.60, 3.3, 0.55,
                 RGBColor(0xFF, 0xFF, 0xFF), C_ORANGE,
                 "メールが迷惑フォルダへ", 10, C_DARK, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide6, s6_left + 3.35, sf_y + 0.68, 0.4, 0.4, "→", 16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left + 3.8, sf_y + 0.60, 3.3, 0.55,
                 RGBColor(0xFF, 0xFF, 0xFF), C_ORANGE,
                 "被害者「メール来てません」", 10, C_DARK, align=PP_ALIGN.CENTER)
add_textbox(slide6, s6_left + 7.15, sf_y + 0.68, 0.4, 0.4, "→", 16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left + 7.6, sf_y + 0.60, 3.5, 0.55,
                 C_HUMAN_BG, C_ORANGE,
                 "犯人「迷惑メールフォルダを\n確認してください」", 10, C_DARK, align=PP_ALIGN.CENTER)
add_textbox(slide6, s6_left + 11.15, sf_y + 0.68, 0.4, 0.4, "→", 16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_bordered_box(slide6, s6_left + 11.6, sf_y + 0.60, 2.9, 0.55,
                 C_DANGER_BG, C_RED,
                 "被害者が自ら発見\n→ 攻撃続行", 10, C_DARK, bold=True, align=PP_ALIGN.CENTER)

# Flaw point
add_rect(slide6, s6_left, sf_y + 1.30, 14.5, 0.50,
         C_RED, "⚠ 破綻ポイント：「迷惑フォルダを確認して」＝ 銀行が絶対に言わないセリフ ＝ 100%詐欺", 13, C_WHITE)

# Worse: reinforces narrative
add_rect(slide6, s6_left, sf_y + 1.90, 14.5, 0.70,
         C_DECISION, "", 1, C_DARK)
add_textbox(slide6, s6_left + 0.2, sf_y + 1.93, 14.0, 0.65,
            "しかし気づけない。さらに悪いことに、迷惑フォルダに入っていたことが手順②の話を補強してしまう\n"
            "「重要なメールがフィルタに引っかかっていた → だから口座更新の案内に気づけなかったんだ」\n"
            "→ 被害者はますます犯人の話を信じる",
            font_size=11, color=C_DARK)

# Bottom conclusion
add_rect(slide6, s6_left, sf_y + 2.80, 14.5, 0.50,
         C_DARK, "手順②の恐怖が全ての破綻ポイント（③メアド、④迷惑フォルダ）を無効化している", 14, C_WHITE)

add_textbox(slide6, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 手順④ フィッシングメールと迷惑メールフィルタ回避",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 7: DNS/DMARCで手順④を技術的に遮断する
# ══════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
slide7.background.fill.solid()
slide7.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide7, 0, 0.2, 16.54, 0.6,
            "技術的対策：DMARC＋DNSで手順④を確実に遮断できる",
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_GREEN)
add_textbox(slide7, 0, 0.75, 16.54, 0.4,
            "顧客の注意力に一切依存せず、銀行側の設定だけでKill Chainを切れる。技術的にはもう解決できる話。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# DMARC policy comparison
s7_left = 0.8
s7_top = 1.3

add_rect(slide7, s7_left, s7_top, 14.9, 0.45,
         C_DARK, "DMARCポリシーの違い — p=reject で正規ドメイン詐称を完全遮断", 13, C_WHITE)

dmarc_rows = [
    ["p=none（監視のみ）",     "偽メールが普通に届く",           "攻撃成立",   C_RED],
    ["p=quarantine（隔離）",   "迷惑フォルダに入る",             "電話で突破される", C_ORANGE],
    ["p=reject（拒否）★",      "メールサーバが受信拒否。配信されない。", "攻撃が止まる",  C_GREEN],
]
dmarc_cols = [3.5, 5.5, 3.5, 2.4]

for ri, row in enumerate(dmarc_rows):
    y = s7_top + 0.45 + ri * 0.50
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x = s7_left
    add_bordered_box(slide7, x, y, dmarc_cols[0], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 10, C_DARK, bold=True)
    x += dmarc_cols[0]
    add_bordered_box(slide7, x, y, dmarc_cols[1], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 10, C_DARK)
    x += dmarc_cols[1]
    add_rect(slide7, x, y, dmarc_cols[2], 0.50,
             row[3], row[2], 10, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Lookalike domain section
lk_y = s7_top + 2.1
add_rect(slide7, s7_left, lk_y, 14.9, 0.45,
         C_DARK, "類似ドメインも含めた包括的DNS対策", 13, C_WHITE)

dns_measures = [
    ["① 類似ドメインの防御的登録",
     "mufg-security.com, bizstation-update.jp 等を銀行が先に取得",
     "犯人が使えるドメインを潰す"],
    ["② 新規登録ドメイン（NRD）ブロック",
     "登録から日が浅いドメインを受信側で高リスク判定・ブロック",
     "犯人が新ドメインを取得しても即座に使えない"],
    ["③ 正規送信ドメインの公開＋GW制御",
     "銀行が「当行のメールは @mufg.jp, @bk.mufg.co.jp のみ」と宣言",
     "企業のメールGWで上記以外を自動ブロック"],
    ["④ SPF/DKIM未設定ドメインの拒否",
     "受信側で認証未設定のドメインからのメールを拒否",
     "犯人の急造ドメインは認証が通らない"],
]

dns_cols = [4.0, 6.5, 4.4]
# Header
add_rect(slide7, s7_left, lk_y + 0.50, dns_cols[0], 0.40,
         C_TAG_BG, "対策", 10, C_WHITE)
add_rect(slide7, s7_left + dns_cols[0], lk_y + 0.50, dns_cols[1], 0.40,
         C_TAG_BG, "内容", 10, C_WHITE)
add_rect(slide7, s7_left + dns_cols[0] + dns_cols[1], lk_y + 0.50, dns_cols[2], 0.40,
         C_TAG_BG, "効果", 10, C_WHITE)

for ri, row in enumerate(dns_measures):
    y = lk_y + 0.90 + ri * 0.50
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x = s7_left
    add_bordered_box(slide7, x, y, dns_cols[0], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 9, C_DARK, bold=True)
    x += dns_cols[0]
    add_bordered_box(slide7, x, y, dns_cols[1], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 9, C_DARK)
    x += dns_cols[1]
    add_bordered_box(slide7, x, y, dns_cols[2], 0.50,
                     bg, C_GREEN, row[2], 9, C_DARK, bold=True)

# Coverage matrix
mx_y = lk_y + 3.1
add_rect(slide7, s7_left, mx_y, 14.9, 0.40,
         C_DARK, "対策カバレッジ", 12, C_WHITE)

mx_headers = ["対策", "正規ドメイン詐称", "類似ドメイン", "新規ドメイン"]
mx_data = [
    ["DMARC p=reject のみ",              "✅ 防げる", "❌ 防げない", "❌ 防げない"],
    ["＋ 類似ドメイン防御的登録",          "✅ 防げる", "✅ 防げる",   "❌ 防げない"],
    ["＋ NRDブロック＋SPF/DKIM必須化",     "✅ 防げる", "✅ 防げる",   "✅ 防げる"],
    ["＋ 正規ドメイン公開＋GW制御（全部）", "✅ 防げる", "✅ 防げる",   "✅ 防げる"],
]
mx_cols = [5.0, 3.3, 3.3, 3.3]

x = s7_left
for ci, h in enumerate(mx_headers):
    add_rect(slide7, x, mx_y + 0.40, mx_cols[ci], 0.40,
             C_TAG_BG, h, 9, C_WHITE)
    x += mx_cols[ci]

for ri, row in enumerate(mx_data):
    y = mx_y + 0.80 + ri * 0.40
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x = s7_left
    add_bordered_box(slide7, x, y, mx_cols[0], 0.40,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 8, C_DARK, bold=True)
    for ci in range(1, 4):
        x += mx_cols[ci - 1]
        is_ok = "✅" in row[ci]
        cell_bg = RGBColor(0xEA, 0xFA, 0xF1) if is_ok else RGBColor(0xFD, 0xED, 0xEC)
        cell_color = C_GREEN if is_ok else C_RED
        add_rect(slide7, x, y, mx_cols[ci], 0.40,
                 cell_bg, row[ci], 9, cell_color, bold=True, align=PP_ALIGN.CENTER)

# Conclusion
add_rect(slide7, s7_left, mx_y + 2.5, 14.9, 0.55,
         C_GREEN,
         "全部やれば手順④のメールは顧客に届かない。Kill Chainはここで切れる。顧客の注意力に一切依存しない。",
         13, C_WHITE)
add_rect(slide7, s7_left, mx_y + 3.1, 14.9, 0.45,
         C_DARK,
         "技術的にはもう解決できる話。やるかやらないかだけ。",
         14, C_WHITE)

add_textbox(slide7, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — DMARC＋DNSによる技術的遮断",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 8: 手順⑤ メールURL — 銀行方針＋金融庁要請に反する
# ══════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
slide8.background.fill.solid()
slide8.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

C_VIOLET = RGBColor(0x8E, 0x44, 0xAD)

add_textbox(slide8, 0, 0.2, 16.54, 0.6,
            "手順⑤「メールのボタンをクリックして」— 最大の破綻ポイント",
            font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_VIOLET)
add_textbox(slide8, 0, 0.75, 16.54, 0.4,
            "銀行の方針・金融庁の要請・顧客への周知、すべてに反する行為。それでも突破される。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

# FSA notice
s8_left = 0.8
s8_top = 1.2
add_rect(slide8, s8_left, s8_top, 14.9, 0.55,
         C_RED,
         "金融庁の要請：金融機関が顧客に送信するメールにURLリンクを記載しないこと",
         13, C_WHITE)
add_textbox(slide8, s8_left, s8_top + 0.58, 14.9, 0.3,
            "→ 多くの銀行が「メール内のリンクはクリックしないでください」と顧客に周知済み",
            9, color=C_GREY_TXT)

# Comparison table
cmp_y = s8_top + 1.0
cmp_cols = [2.5, 6.0, 6.4]
cmp_headers = ["", "銀行の正規誘導", "この攻撃"]
cmp_rows = [
    ["方法",       "ブックマーク・URL直接入力・公式アプリ",   "メールのボタンをクリック"],
    ["URLの出所",  "顧客が自分で持っている",                 "犯人が送りつけた"],
    ["顧客の行動", "能動的（自分でアクセス）",               "受動的（指示されてクリック）"],
    ["利便性",     "面倒（ブックマーク探す、URL入力）",       "楽（ボタン1つ押すだけ）"],
]

x = s8_left
for ci, h in enumerate(cmp_headers):
    add_rect(slide8, x, cmp_y, cmp_cols[ci], 0.45,
             C_DARK, h, 10, C_WHITE)
    x += cmp_cols[ci]

for ri, row in enumerate(cmp_rows):
    y = cmp_y + 0.45 + ri * 0.50
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    x = s8_left
    add_bordered_box(slide8, x, y, cmp_cols[0], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 10, C_DARK, bold=True)
    x += cmp_cols[0]
    add_bordered_box(slide8, x, y, cmp_cols[1], 0.50,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 10, C_DARK)
    x += cmp_cols[1]
    add_bordered_box(slide8, x, y, cmp_cols[2], 0.50,
                     bg, C_RED, row[2], 10, C_RED, bold=True)

# Irony section
irn_y = cmp_y + 2.55
add_rect(slide8, s8_left, irn_y, 14.9, 0.90,
         C_VIOLET, "", 1, C_WHITE)
add_textbox(slide8, s8_left + 0.3, irn_y + 0.05, 14.3, 0.85,
            "皮肉：犯人のほうが「ユーザー体験が良い」\n"
            "正規の銀行：「ブックマークからアクセスしてください」（面倒）\n"
            "犯人：「ボタンをクリックするだけです」（簡単）\n"
            "→ 焦っている被害者がどちらに従いたくなるかは明白",
            font_size=11, color=C_WHITE, bold=True)

# Flaw summary
fw_y = irn_y + 1.10
add_rect(slide8, s8_left, fw_y, 14.9, 0.40,
         C_DARK, "手順⑤の破綻ポイントまとめ", 12, C_WHITE)

fw_rows = [
    ["銀行の方針",      "メールにURLリンクを載せない"],
    ["金融庁の要請",    "金融機関のメールにURL記載禁止"],
    ["顧客への周知",    "「メールのリンクをクリックしないで」と案内済み"],
]
for ri, row in enumerate(fw_rows):
    y = fw_y + 0.40 + ri * 0.40
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    add_bordered_box(slide8, s8_left, y, 4.0, 0.40,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 10, C_DARK, bold=True)
    add_bordered_box(slide8, s8_left + 4.0, y, 10.9, 0.40,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 10, C_DARK)

# "Still breaks through" row
y = fw_y + 0.40 + 3 * 0.40
add_rect(slide8, s8_left, y, 4.0, 0.40,
         C_RED, "それでも突破される理由", 10, C_WHITE, bold=True)
add_rect(slide8, s8_left + 4.0, y, 10.9, 0.40,
         C_RED, "手順②の恐怖 ＋ 電話でのリアルタイム指示 ＋ 犯人のほうがUXが良い", 10, C_WHITE)

# Bottom conclusion
add_rect(slide8, s8_left, fw_y + 2.05, 14.9, 0.55,
         C_DARK,
         "注意喚起・ガイドライン・金融庁の要請 — すべてが「冷静な時の判断力」を前提にしている。この攻撃は、その前提自体を手順②で破壊している。",
         12, C_WHITE)

add_textbox(slide8, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — 手順⑤ メールURLと金融庁要請",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# Slide 9: 「ブックマークからアクセスして」は機能しない
# ══════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
slide9.background.fill.solid()
slide9.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

add_textbox(slide9, 0, 0.2, 16.54, 0.6,
            "「ブックマークからアクセスして」は機能しない",
            font_size=24, bold=True, align=PP_ALIGN.CENTER, color=C_ORANGE)
add_textbox(slide9, 0, 0.75, 16.54, 0.4,
            "ブックマーク登録を顧客に任せている時点で、対策として成立していない。",
            font_size=12, align=PP_ALIGN.CENTER, color=C_GREY_TXT)

s9_left = 0.8
s9_top = 1.3

# How bank conveys URL
add_rect(slide9, s9_left, s9_top, 7.2, 0.45,
         C_DARK, "銀行のURL伝達方法（法人IB）", 12, C_WHITE)
url_methods = [
    ["① 契約時の書面", "紙にURLが記載されている"],
    ["② 銀行の窓口", "担当者から口頭で案内"],
    ["③ 銀行の公式HP", "検索 → HP → IBログインページ"],
]
for ri, row in enumerate(url_methods):
    y = s9_top + 0.45 + ri * 0.45
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    add_bordered_box(slide9, s9_left, y, 2.5, 0.45,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[0], 10, C_DARK, bold=True)
    add_bordered_box(slide9, s9_left + 2.5, y, 4.7, 0.45,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 10, C_DARK)

add_textbox(slide9, s9_left, s9_top + 1.85, 7.2, 0.3,
            "→ 顧客が自分でブックマーク登録する必要がある", 10, bold=True, color=C_RED)

# Reality vs assumption
add_rect(slide9, s9_left + 7.7, s9_top, 7.0, 0.45,
         C_DARK, "銀行の想定 vs 現実", 12, C_WHITE)
reality_rows = [
    ["契約時の書面を保管", "どこにしまったかわからない"],
    ["ブックマークに登録済み", "登録していない人が多い"],
    ["毎回ブックマークからアクセス", "毎回「BizSTATION」で検索している"],
]
for ri, row in enumerate(reality_rows):
    y = s9_top + 0.45 + ri * 0.45
    add_bordered_box(slide9, s9_left + 7.7, y, 3.2, 0.45,
                     RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xCC, 0xCC, 0xCC),
                     row[0], 9, C_DARK)
    add_rect(slide9, s9_left + 10.9, y, 3.8, 0.45,
             C_DANGER_BG, row[1], 9, C_RED, bold=True, align=PP_ALIGN.LEFT)

# Impact on attack
imp_y = s9_top + 2.4
add_rect(slide9, s9_left, imp_y, 14.9, 0.45,
         C_DARK, "ブックマーク登録の有無が攻撃成功率を左右する", 12, C_WHITE)

# Two columns
add_rect(slide9, s9_left, imp_y + 0.50, 7.2, 0.40,
         C_GREEN, "ブックマーク登録済みの人（少数）", 10, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide9, s9_left, imp_y + 0.90, 7.2, 0.70,
                 RGBColor(0xFF, 0xFF, 0xFF), C_GREEN,
                 "→ 「自分のブックマークと違うURLだ」と気づける可能性がある\n→ でも電話中なので確認しない",
                 10, C_DARK)

add_rect(slide9, s9_left + 7.7, imp_y + 0.50, 7.0, 0.40,
         C_RED, "ブックマーク未登録の人（多数）", 10, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bordered_box(slide9, s9_left + 7.7, imp_y + 0.90, 7.0, 0.70,
                 RGBColor(0xFF, 0xFF, 0xFF), C_RED,
                 "→ 普段から検索やメールリンクでアクセスしている\n→ メールのボタンをクリックすることに抵抗がない\n→ 攻撃にそのまま乗ってしまう",
                 10, C_DARK)

# What bank should have done
shd_y = imp_y + 1.85
add_rect(slide9, s9_left, shd_y, 14.9, 0.45,
         C_GREEN, "銀行がやるべきだったこと", 12, C_WHITE)

should_rows = [
    ["① 契約時にブックマークを設定してあげる", "窓口対応 or 初期設定ツールで銀行側が登録"],
    ["② 公式アプリでしかアクセスできない仕組み", "ブラウザ経由のアクセスを廃止"],
    ["③ 専用クライアントに統一", "ブラウザを使わせないことでフィッシング自体を不可能にする"],
]
for ri, row in enumerate(should_rows):
    y = shd_y + 0.45 + ri * 0.45
    bg = RGBColor(0xF9, 0xF9, 0xF9) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    add_bordered_box(slide9, s9_left, y, 5.5, 0.45,
                     bg, C_GREEN, row[0], 9, C_DARK, bold=True)
    add_bordered_box(slide9, s9_left + 5.5, y, 9.4, 0.45,
                     bg, RGBColor(0xCC, 0xCC, 0xCC), row[1], 9, C_DARK)

# Conclusion
add_rect(slide9, s9_left, shd_y + 1.85, 14.9, 0.50,
         C_DARK,
         "「自分でブックマーク登録してね」では対策にならない — 銀行側がアクセス経路を制御すべき",
         13, C_WHITE)

add_textbox(slide9, 10, 11.35, 6.5, 0.25,
            "Fraud Kill Chain — ブックマーク対策の限界",
            7, color=C_GREY_TXT, align=PP_ALIGN.RIGHT)

# ── Save ──
out = "/home/user/security-news-pwa/Fraud_Kill_Chain.pptx"
prs.save(out)
print(f"Saved: {out}")
